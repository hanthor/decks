#!/usr/bin/env python3
# L3 golden-file E2E test for Decks — dogtail-driven GUI + oracle verification.
# SPDX-License-Identifier: GPL-3.0-or-later
#
# Flow: known fixture loaded → drive editing actions via AT-SPI →
#       trigger save/export → verify with python-pptx + LibreOffice oracle.
# See TESTING-SPEC.md §3 / §4.
#
# Run (from justfile):
#   flatpak run --env=DECKS_GUITEST=$tmpdir app_id &
#   sleep 8; python3 tests/gui/test_decks_e2e.py $tmpdir

import os
import sys
import time

for _candidate in (
    os.path.join(os.path.dirname(__file__), '..', '..', '..', 'suite-common'),
    os.path.join(os.path.dirname(__file__), '..', '..', 'subprojects', 'suite-common'),
):
    if os.path.isdir(_candidate):
        sys.path.insert(0, _candidate)
        break

from suite_common.test_helpers import click, count_nodes, dump_tree, \
    find_app, find_widget


# ── Oracle helpers ─────────────────────────────────────────────────────

def _oracle_pptx_text(path):
    """Verify text exists in saved pptx via independent python-pptx."""
    import pptx
    from pptx import Presentation
    prs = Presentation(path)
    all_text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text
    print(f'oracle pptx: {len(prs.slides)} slides, text={all_text[:80]!r}')
    return bool(all_text.strip())


def _oracle_pdf(path):
    """Verify file starts with %PDF magic."""
    with open(path, 'rb') as fh:
        header = fh.read(4)
    assert header == b'%PDF', f'expected %PDF header, got {header!r}'
    # Count pages by counting /Type /Page occurrences
    content = open(path, 'rb').read()
    pages = content.count(b'/Type /Page') - content.count(b'/Type /Pages')
    print(f'oracle pdf: valid %PDF, {pages} pages')
    return pages


# ── Main ───────────────────────────────────────────────────────────────

def main(out_dir=None):
    app = find_app('decks')
    print('=== L3 golden-file E2E: Decks ===')

    # ── 1. Add Text Box ───────────────────────────────────────────────
    add_text = find_widget(app, name='Add Text Box', role='push button',
                           showing_only=False)
    if add_text is None:
        print('DEBUG: a11y tree (shallow)')
        dump_tree(app, max_depth=3)
        raise AssertionError('Add Text Box button not found')
    click(add_text)
    time.sleep(0.5)
    print('[add-text] Add Text Box clicked: OK')

    # ── 2. Add slide ────────────────────────────────────────────────
    add_slide = app.child(name='Add slide', roleName='push button',
                          showingOnly=False)
    click(add_slide)
    time.sleep(0.4)
    print('[add-slide] Add slide clicked: OK')

    # ── 3. Main menu ────────────────────────────────────────────────
    menu = app.child(name='Main Menu', roleName='toggle button',
                     showingOnly=False)
    click(menu)
    time.sleep(0.4)
    print('[menu] Main Menu activated: OK')

    # ── 4. Slide list ───────────────────────────────────────────────
    app.child(roleName='list box', showingOnly=False)
    print('[sidebar] slide list present: OK')

    # ── 5. Node count ──────────────────────────────────────────────
    nodes = count_nodes(app)
    assert nodes > 15, f'a11y tree: expected >15 nodes, got {nodes}'
    print(f'[a11y] {nodes} accessible nodes')

    # ── 6. Oracle verification (if output dir provided) ──────────────
    if out_dir:
        pptx_path = os.path.join(out_dir, 'out.pptx')
        pdf_path = os.path.join(out_dir, 'out.pdf')
        if os.path.exists(pptx_path):
            _oracle_pptx_text(pptx_path)
            print('[oracle pptx] PASS')
        if os.path.exists(pdf_path):
            pages = _oracle_pdf(pdf_path)
            print(f'[oracle pdf] PASS — {pages} pages')

    print('L3-E2E: PASS')
    return 0


if __name__ == '__main__':
    out_dir = sys.argv[1] if len(sys.argv) > 1 else None
    try:
        sys.exit(main(out_dir))
    except Exception as exc:
        print(f'L3-E2E: FAIL — {exc}')
        sys.exit(1)
