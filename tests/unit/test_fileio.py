# L1 adapter round-trip tests for Decks — pure pytest, no display.
# SPDX-License-Identifier: GPL-3.0-or-later
#
# Verifies pptx/odp text, images, and slide count via independent read-back.
# See TESTING-SPEC.md §1.
#
# Run: pytest tests/unit/

import os
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'src'))
import fileio  # noqa: E402


# ── Helpers ────────────────────────────────────────────────────────────

def _slide(text='HelloDeck', left=80, top=120, width=400, height=60,
           font_size=32):
    """Build a minimal single-textbox slide dict."""
    return {
        'version': '5.5.2', 'background': '#ffffff',
        'objects': [
            {'type': 'i-text', 'text': text, 'left': left, 'top': top,
             'width': width, 'height': height, 'fontSize': font_size,
             'fontFamily': 'sans-serif', 'fill': '#202020'},
        ],
    }


def _roundtrip(slides, ext):
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, f'deck.{ext}')
        fileio.write_deck(path, slides)
        return fileio.read_deck(path)


def _independent_pptx_read(path):
    """Read a pptx with python-pptx (independent path from fileio)."""
    pytest.importorskip('pptx')
    from pptx import Presentation
    prs = Presentation(path)
    result = {'slide_count': len(prs.slides), 'slides': []}
    for slide in prs.slides:
        info = {'texts': [], 'pictures': 0}
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                info['texts'].append(shape.text_frame.text)
            if hasattr(shape, 'image'):
                info['pictures'] += 1
        result['slides'].append(info)
    return result


def _independent_odp_read(path):
    """Read an odp with odfpy (independent path from fileio)."""
    pytest.importorskip('odf.opendocument')
    from odf.opendocument import load
    from odf.draw import Page, Frame
    from odf.text import P
    doc = load(path)
    result = {'slide_count': 0, 'slides': []}
    for page in doc.getElementsByType(Page):
        texts = []
        for frame in page.getElementsByType(Frame):
            t = ''.join(str(p) for p in frame.getElementsByType(P))
            if t.strip():
                texts.append(t.strip())
        result['slides'].append({'texts': texts})
        result['slide_count'] += 1
    return result


# ── PPTX text ──────────────────────────────────────────────────────────

class TestPptxText:
    @classmethod
    def setup_class(cls):
        pytest.importorskip('pptx')

    def test_text_roundtrip(self):
        slides = [_slide('HelloDeck'), _slide('Slide Two')]
        result = _roundtrip(slides, 'pptx')
        assert len(result) == 2
        texts = [o.get('text') for s in result
                 for o in s.get('objects', []) if o.get('type') == 'i-text']
        assert 'HelloDeck' in texts
        assert 'Slide Two' in texts

    def test_slide_count_preserved(self):
        slides = [_slide('A'), _slide('B'), _slide('C')]
        result = _roundtrip(slides, 'pptx')
        assert len(result) == 3

    def test_independent_reader_text(self):
        slides = [_slide('HelloDeck')]
        with tempfile.TemporaryDirectory() as td:
            path = os.path.join(td, 'out.pptx')
            fileio.write_deck(path, slides)
            proof = _independent_pptx_read(path)
        assert proof['slide_count'] == 1
        assert 'HelloDeck' in proof['slides'][0]['texts']


# ── PPTX images ────────────────────────────────────────────────────────

IMG_URL = (
    'data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1'
    'HAwCAAAAC0lEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=='
)


def _image_slide():
    return {
        'version': '5.5.2', 'background': '#ffffff',
        'objects': [
            {'type': 'image', 'src': IMG_URL, 'left': 100, 'top': 100,
             'width': 200, 'height': 150, 'scaleX': 1, 'scaleY': 1},
        ],
    }


class TestPptxImages:
    @classmethod
    def setup_class(cls):
        pytest.importorskip('pptx')

    def test_image_roundtrip(self):
        slides = [_image_slide()]
        result = _roundtrip(slides, 'pptx')
        assert len(result) == 1
        obj_types = [o.get('type') for o in result[0].get('objects', [])]
        assert 'image' in obj_types

    def test_independent_reader_image_count(self):
        slides = [_image_slide()]
        with tempfile.TemporaryDirectory() as td:
            path = os.path.join(td, 'out.pptx')
            fileio.write_deck(path, slides)
            proof = _independent_pptx_read(path)
        assert proof['slides'][0]['pictures'] == 1


# ── ODP ────────────────────────────────────────────────────────────────

class TestOdpText:
    @classmethod
    def setup_class(cls):
        pytest.importorskip('odf.opendocument')

    def test_text_roundtrip(self):
        slides = [_slide('HelloOdp')]
        result = _roundtrip(slides, 'odp')
        assert len(result) >= 1
        texts = [o.get('text') for s in result
                 for o in s.get('objects', []) if o.get('type') == 'i-text']
        assert 'HelloOdp' in texts

    def test_slide_count(self):
        slides = [_slide('One'), _slide('Two'), _slide('Three')]
        result = _roundtrip(slides, 'odp')
        assert len(result) >= 3

    def test_independent_reader_text(self):
        slides = [_slide('HelloOdp'), _slide('SlideB')]
        with tempfile.TemporaryDirectory() as td:
            path = os.path.join(td, 'out.odp')
            fileio.write_deck(path, slides)
            proof = _independent_odp_read(path)
        assert proof['slide_count'] == 2
        all_text = ' '.join(t for s in proof['slides'] for t in s['texts'])
        assert 'HelloOdp' in all_text


# ── Empty / edge cases ─────────────────────────────────────────────────

class TestEdgeCases:
    def test_empty_deck_pptx(self):
        pytest.importorskip('pptx')
        result = _roundtrip([], 'pptx')
        assert isinstance(result, list)

    def test_empty_slide_pptx(self):
        pytest.importorskip('pptx')
        slides = [{'version': '5.5.2', 'objects': [], 'background': '#ffffff'}]
        result = _roundtrip(slides, 'pptx')
        assert len(result) >= 1

    def test_empty_deck_odp(self):
        pytest.importorskip('odf.opendocument')
        result = _roundtrip([], 'odp')
        assert isinstance(result, list)
